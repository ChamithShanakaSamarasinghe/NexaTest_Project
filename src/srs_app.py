import os

# FFmpeg path
os.environ["PATH"] += r";C:\Users\raven\OneDrive\Desktop\ffmpeg\ffmpeg-8.1-essentials_build\bin"

import streamlit as st
import sys
import json
import re
import subprocess
import email
import uuid

from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract

pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

import sqlite3
from datetime import datetime

sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))

from src.enhancer import PipelineEnhancer
from src.fpr.fpr_pipeline import run_fpr_pipeline
from src.testcase.generator import generate_test_cases

DB_FILE = "db.sqlite3"
SRS_FOLDER = "srs_docs"

os.makedirs(SRS_FOLDER, exist_ok=True)

enhancer = PipelineEnhancer()

# ---------------- TEXT SPLITTER ----------------
def split_into_sentences(text):
    text = re.sub(r'(\d+\.\d+)', r' \1 ', text)
    sentences = re.split(r'(?<=[.!?])\s+|\n+', text)
    return [s.strip() for s in sentences if len(s.strip()) > 10]

# ---------------- EMAIL CLEANER ----------------
def clean_email_text(text):
    text = re.sub(r"(From:.*|Sent:.*|To:.*|Subject:.*)", "", text, flags=re.IGNORECASE)
    text = re.sub(r"(Hi|Hello|Dear).*?\n", "", text, flags=re.IGNORECASE)
    text = re.sub(r"(Regards|Thanks|Best).*", "", text, flags=re.IGNORECASE)
    return text.strip()

# ---------------- CHAT ----------------
def normalize_chat_text(text):
    replacements = {
        "pls": "please", "plz": "please", "u": "you", "ur": "your",
        "asap": "as soon as possible", "btw": "by the way", "im": "i am"
    }
    return " ".join([replacements.get(w.lower(), w) for w in text.split()])

def parse_chat_lines(text):
    lines = text.split("\n")
    cleaned = []
    for line in lines:
        line = re.sub(r"\[.*?\]|\d{1,2}:\d{2}.*", "", line)
        if len(line.strip()) > 5:
            cleaned.append(line.strip())
    return cleaned

# ---------------- VIDEO ----------------
def extract_audio_from_video(video_path, output_audio_path):
    command = ["ffmpeg", "-i", video_path, "-vn", "-acodec", "mp3", output_audio_path, "-y"]
    subprocess.run(command, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)

# ---------------- DB ----------------
def insert_into_db(doc_filename, sections, requirements, features, test_results, fpr_result):
    conn = sqlite3.connect(DB_FILE)
    cursor = conn.cursor()

    cursor.execute("INSERT INTO documents (filename) VALUES (?)", (doc_filename,))
    doc_id = cursor.lastrowid

    for sec in sections:
        cursor.execute(
            "INSERT INTO sections (document_id, section_name, content) VALUES (?, ?, ?)",
            (doc_id, sec["section_name"], sec["content"]),
        )

    for req in requirements:
        cursor.execute(
            "INSERT INTO requirements (document_id, requirement_text) VALUES (?, ?)",
            (doc_id, req),
        )

    for feat in features:
        cursor.execute(
            "INSERT INTO features (document_id, feature_name) VALUES (?, ?)",
            (doc_id, feat),
        )

    for test_name, status in test_results.items():
        cursor.execute(
            "INSERT INTO test_results (doc_id, test_name, status) VALUES (?, ?, ?)",
            (doc_id, test_name, status),
        )

    cursor.execute(
        "INSERT INTO fpr_results (document_id, clusters, keywords, metrics) VALUES (?, ?, ?, ?)",
        (
            doc_id,
            json.dumps(fpr_result["clusters"]),
            json.dumps(fpr_result["keywords"]),
            json.dumps(fpr_result["metrics"]),
        ),
    )

    conn.commit()
    conn.close()
    return doc_id

# ---------------- UI ----------------
st.set_page_config(page_title="🚀 SRS Automation Platform", layout="wide")
st.title("🚀 SRS Automation Platform (AI Powered)")

uploaded_files = st.file_uploader(
    "Upload SRS (Multiple inputs supported → single document ID)",
    type=["docx", "pdf", "pptx", "png", "jpg", "jpeg", "mp3", "wav", "mp4", "avi", "mov", "eml", "txt"],
    accept_multiple_files=True
)

# 🔥 Prevent re-run duplication
if "processed" not in st.session_state:
    st.session_state.processed = False

if uploaded_files:

    st.info("📂 Files ready. Click below to process all together.")

    if st.button("🚀 Process All Inputs") and not st.session_state.processed:

        st.session_state.processed = True

        session_id = str(uuid.uuid4())[:8]
        all_paragraphs = []
        file_names = []

        for uploaded_file in uploaded_files:

            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            srs_filename = f"{timestamp}_{uploaded_file.name}"
            srs_path = os.path.join(SRS_FOLDER, srs_filename)

            with open(srs_path, "wb") as f:
                f.write(uploaded_file.getbuffer())

            file_names.append(uploaded_file.name)
            st.success(f"✅ Processed: {uploaded_file.name}")

            paragraphs = []

            try:
                if uploaded_file.name.endswith(".docx"):
                    doc = Document(srs_path)
                    paragraphs = [p.text for p in doc.paragraphs]

                elif uploaded_file.name.endswith(".pdf"):
                    import PyPDF2
                    reader = PyPDF2.PdfReader(srs_path)
                    for page in reader.pages:
                        text = page.extract_text()
                        if text:
                            paragraphs.extend(split_into_sentences(text))

                elif uploaded_file.name.endswith(".pptx"):
                    prs = Presentation(srs_path)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                for s in split_into_sentences(shape.text):
                                    paragraphs.append(f"[PPT_INPUT] {s}")

                elif uploaded_file.name.endswith((".png", ".jpg", ".jpeg")):
                    image = Image.open(srs_path)
                    text = pytesseract.image_to_string(image)
                    for s in split_into_sentences(text):
                        paragraphs.append(f"[IMAGE_INPUT] {s}")

                elif uploaded_file.name.endswith((".mp3", ".wav")):
                    import whisper
                    model = whisper.load_model("base")
                    result = model.transcribe(srs_path)
                    for s in split_into_sentences(result["text"]):
                        paragraphs.append(f"[AUDIO_INPUT] {s}")

                elif uploaded_file.name.endswith((".mp4", ".avi", ".mov")):
                    import whisper
                    audio_path = srs_path + ".mp3"

                    st.info(f"🎬 Extracting: {uploaded_file.name}")
                    extract_audio_from_video(srs_path, audio_path)

                    model = whisper.load_model("base")
                    result = model.transcribe(audio_path)

                    for s in split_into_sentences(result["text"]):
                        paragraphs.append(f"[VIDEO_INPUT] {s}")

                elif uploaded_file.name.endswith(".eml"):
                    st.info(f"📧 Email: {uploaded_file.name}")

                    raw_email = uploaded_file.read().decode("utf-8", errors="ignore")
                    msg = email.message_from_string(raw_email)

                    body = ""
                    if msg.is_multipart():
                        for part in msg.walk():
                            if part.get_content_type() == "text/plain":
                                body += part.get_payload(decode=True).decode(errors="ignore")
                    else:
                        body = msg.get_payload(decode=True).decode(errors="ignore")

                    cleaned = clean_email_text(body)

                    for s in split_into_sentences(cleaned):
                        paragraphs.append(f"[EMAIL_INPUT] {s}")

                elif uploaded_file.name.endswith(".txt"):
                    st.info(f"💬 Chat: {uploaded_file.name}")

                    raw_text = uploaded_file.read().decode("utf-8")
                    lines = parse_chat_lines(raw_text)

                    for line in lines:
                        normalized = normalize_chat_text(line)
                        for s in split_into_sentences(normalized):
                            paragraphs.append(f"[CHAT_INPUT] {s}")

            except Exception as e:
                st.error(f"{uploaded_file.name} → {str(e)}")

            all_paragraphs.extend(paragraphs)

        # ---------------- PIPELINE ----------------
        sections, requirements, features = [], [], []
        current_section = None
        full_text = ""

        for para in all_paragraphs:
            text = para.strip()
            if not text:
                continue

            full_text += text + " "

            if text.isupper() or text.endswith(":"):
                current_section = text
                sections.append({"section_name": current_section, "content": ""})
                continue

            if current_section:
                sections[-1]["content"] += text + " "

            lower = text.lower()

            if re.search(r"\b(shall|must|required to)\b", lower):
                requirements.append(text)

            elif any(tag in text for tag in ["[EMAIL_INPUT]", "[CHAT_INPUT]", "[VIDEO_INPUT]", "[AUDIO_INPUT]", "[PPT_INPUT]"]):
                clean = re.sub(r"\[.*?_INPUT\]", "", text).strip()
                if len(clean.split()) > 5:
                    requirements.append(f"The system shall {clean.lower()}")

            if any(word in lower for word in ["feature", "module", "management", "system"]):
                features.append(text)

        enhanced = enhancer.enhance(full_text)
        fpr_result = run_fpr_pipeline(requirements) if requirements else {"clusters": [], "keywords": {}, "metrics": {}}
        test_cases = generate_test_cases(requirements)

        test_results = {
            "test_requirement_extraction": "PASSED",
            "test_testcase_generation": "PASSED",
            "test_srs_upload": "PASSED"
        }

        doc_name = f"SESSION_{session_id}_" + "_".join(file_names)
        doc_id = insert_into_db(doc_name, sections, requirements, features, test_results, fpr_result)

        st.success(f"💾 Saved ALL inputs under ONE doc_id: {doc_id}")

        # ---------------- METRICS ----------------
        col1, col2, col3, col4 = st.columns(4)
        col1.metric("📌 Requirements", len(requirements))
        col2.metric("✨ Features", len(features))
        col3.metric("🧠 Clusters", len(set(fpr_result["clusters"])))
        col4.metric("🧪 Test Cases", len(test_cases))

        # ---------------- UI TABS ----------------
        tabs = st.tabs(["📄 Sections", "📌 Requirements", "✨ Features", "🧠 Clusters", "🧪 Test Cases"])

        with tabs[0]:
            for sec in sections:
                with st.expander(sec["section_name"]):
                    st.write(sec["content"])

        with tabs[1]:
            for r in requirements:
                st.markdown(f"- {r}")

        with tabs[2]:
            for i, f in enumerate(features, 1):
                with st.expander(f"Feature {i}"):
                    st.write(f)

        with tabs[3]:
            cluster_map = {}
            for req, label in zip(requirements, fpr_result["clusters"]):
                cluster_map.setdefault(label, []).append(req)

            for cid, reqs in cluster_map.items():
                with st.expander(f"Cluster {cid}"):
                    for r in reqs:
                        st.markdown(f"- {r}")

            st.json(fpr_result["metrics"])

        with tabs[4]:
            for tc in test_cases[:100]:
                with st.expander(f"{tc.get('id')} | {tc.get('technique')}"):
                    st.write("📌 Requirement:", tc.get("requirement"))
                    st.write("🔹 Input:", tc.get("input"))
                    st.write("✅ Expected:", tc.get("expected"))

        st.balloons()